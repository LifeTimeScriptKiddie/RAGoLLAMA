"""
Document Processor

Handles extraction and chunking of content from various document formats.
"""

from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path
import logging
import hashlib
import magic
from PIL import Image
import pytesseract

# Document format processors
import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Processes various document formats and extracts text content."""
    
    def __init__(self, config: Dict[str, Any]):
        """
        Initialize document processor.
        
        Args:
            config: Configuration dictionary
        """
        self.config = config
        self.max_chunk_size = config.get("max_chunk_size", 1000)
        self.chunk_overlap = config.get("chunk_overlap", 200)
        self.supported_formats = config.get("supported_formats", [".pdf", ".txt", ".docx", ".pptx", ".jpg", ".png"])
        
    def process_document(self, file_path: str) -> Tuple[List[str], Dict[str, Any]]:
        """
        Process a document and extract text chunks.
        
        Args:
            file_path: Path to the document file
            
        Returns:
            Tuple of (text chunks, metadata)
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"Document not found: {file_path}")
        
        # Detect file type
        file_type = self._detect_file_type(file_path)
        
        if file_type not in [fmt.lstrip('.') for fmt in self.supported_formats]:
            raise ValueError(f"Unsupported file format: {file_type}")
        
        logger.info(f"Processing {file_type} document: {file_path.name}")
        
        # Extract text based on file type
        if file_type == "pdf":
            text = self._extract_pdf_text(file_path)
        elif file_type == "txt":
            text = self._extract_txt_text(file_path)
        elif file_type == "docx":
            text = self._extract_docx_text(file_path)
        elif file_type == "pptx":
            text = self._extract_pptx_text(file_path)
        elif file_type in ["jpg", "jpeg", "png"]:
            text = self._extract_image_text(file_path)
        else:
            raise ValueError(f"Handler not implemented for: {file_type}")
        
        # Create document chunks
        chunks = self._create_chunks(text)
        
        # Generate metadata
        metadata = self._create_metadata(file_path, file_type, text, chunks)
        
        logger.info(f"Extracted {len(chunks)} chunks from {file_path.name}")
        return chunks, metadata
    
    def _detect_file_type(self, file_path: Path) -> str:
        """Detect file type using magic numbers and extension."""
        try:
            mime_type = magic.from_file(str(file_path), mime=True)
            
            # Map MIME types to our format names
            mime_map = {
                "application/pdf": "pdf",
                "text/plain": "txt",
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
                "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
                "image/jpeg": "jpg",
                "image/png": "png"
            }
            
            detected_type = mime_map.get(mime_type)
            if detected_type:
                return detected_type
                
        except Exception as e:
            logger.warning(f"Magic detection failed: {e}")
        
        # Fallback to extension
        extension = file_path.suffix.lower().lstrip(".")
        return extension
    
    def _extract_pdf_text(self, file_path: Path) -> str:
        """Extract text from PDF file."""
        try:
            text = ""
            with open(file_path, "rb") as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
            return text.strip()
        except Exception as e:
            logger.error(f"Error extracting PDF text: {e}")
            raise
    
    def _extract_txt_text(self, file_path: Path) -> str:
        """Extract text from TXT file."""
        try:
            encodings = ["utf-8", "latin-1", "cp1252"]
            for encoding in encodings:
                try:
                    with open(file_path, "r", encoding=encoding) as file:
                        return file.read()
                except UnicodeDecodeError:
                    continue
            raise ValueError("Unable to decode text file with common encodings")
        except Exception as e:
            logger.error(f"Error extracting TXT text: {e}")
            raise
    
    def _extract_docx_text(self, file_path: Path) -> str:
        """Extract text from DOCX file."""
        try:
            doc = DocxDocument(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text.strip()
        except Exception as e:
            logger.error(f"Error extracting DOCX text: {e}")
            raise
    
    def _extract_pptx_text(self, file_path: Path) -> str:
        """Extract text from PPTX file."""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip()
        except Exception as e:
            logger.error(f"Error extracting PPTX text: {e}")
            raise
    
    def _extract_image_text(self, file_path: Path) -> str:
        """Extract text from image using OCR."""
        try:
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)
            return text.strip()
        except Exception as e:
            logger.error(f"Error extracting image text with OCR: {e}")
            # Return empty string if OCR fails
            return ""
    
    def _create_chunks(self, text: str) -> List[str]:
        """Split text into overlapping chunks."""
        if not text:
            return []
        
        chunks = []
        start = 0
        text_length = len(text)
        
        while start < text_length:
            end = start + self.max_chunk_size
            
            # If this is not the last chunk, try to break at word boundary
            if end < text_length:
                # Look for the last space before the end position
                last_space = text.rfind(" ", start, end)
                if last_space > start:
                    end = last_space
            
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            # Move start position with overlap
            start = end - self.chunk_overlap
            if start <= 0:
                start = end
        
        return chunks
    
    def _create_metadata(self, file_path: Path, file_type: str, text: str, chunks: List[str]) -> Dict[str, Any]:
        """Create metadata for the processed document."""
        # Calculate file hash
        with open(file_path, "rb") as f:
            file_hash = hashlib.md5(f.read()).hexdigest()
        
        metadata = {
            "file_path": str(file_path),
            "file_name": file_path.name,
            "file_type": file_type,
            "file_size": file_path.stat().st_size,
            "file_hash": file_hash,
            "text_length": len(text),
            "num_chunks": len(chunks),
            "max_chunk_size": self.max_chunk_size,
            "chunk_overlap": self.chunk_overlap,
            "chunk_lengths": [len(chunk) for chunk in chunks],
            "processed_at": __import__("datetime").datetime.utcnow().isoformat()
        }
        
        return metadata
    
    def get_supported_formats(self) -> List[str]:
        """Get list of supported document formats."""
        return self.supported_formats